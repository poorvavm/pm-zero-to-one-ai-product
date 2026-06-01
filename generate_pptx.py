#!/usr/bin/env python3
"""Generate PharCovAI presentation PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette (PharCovAI: Light Blue, Cyan, White) ──
CYAN        = RGBColor(0x0D, 0x94, 0x88)  # #0D9488 teal/cyan accent
LIGHT_BLUE  = RGBColor(0x0E, 0xA5, 0xE9)  # #0EA5E9
DARK_BLUE   = RGBColor(0x0F, 0x17, 0x2A)  # dark navy for text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF8, 0xF9, 0xFA)
MID_GRAY    = RGBColor(0x6B, 0x72, 0x80)
BORDER_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
RED         = RGBColor(0xDC, 0x26, 0x26)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)
AMBER       = RGBColor(0xCA, 0x8A, 0x04)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_slide_gradient(slide, color1, color2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=DARK_BLUE, alignment=PP_ALIGN.LEFT,
                      font_name="Calibri", bold=False, line_spacing=1.3,
                      bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet:
            p.text = f"  {line}"
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow(slide, left, top, width, height, color=MID_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, text="",
               font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_table(slide, left, top, width, height, rows, cols,
              header_data, row_data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, hdr in enumerate(header_data):
        cell = table.cell(0, j)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = CYAN

    for i, row in enumerate(row_data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_BLUE
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY

    return table_shape


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_stat_card(slide, left, top, number, label, color=CYAN):
    card_w = Inches(2.5)
    card_h = Inches(2.0)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER_GRAY
    shape.line.width = Pt(1)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.3),
                 card_w - Inches(0.4), Inches(0.9),
                 number, font_size=36, bold=True, color=color,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.2),
                 card_w - Inches(0.4), Inches(0.6),
                 label, font_size=13, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)


def add_section_number(slide, number):
    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.45),
                 f"{number:02d}", font_size=14, color=CYAN, bold=True)


def add_slide_title(slide, title, subtitle=None):
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.7),
                 title, font_size=32, bold=True, color=DARK_BLUE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.05), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=MID_GRAY)


def add_bottom_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, H - Inches(0.08), W, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()


def add_divider_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER_GRAY
    shape.line.fill.background()


# ── SLIDE 1: Title ───────────────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)

# Gradient accent bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.15))
bar.fill.solid()
bar.fill.fore_color.rgb = CYAN
bar.line.fill.background()

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "PharCovAI", font_size=54, bold=True, color=DARK_BLUE)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "AI-Powered Pharmacovigilance & Compliance Intelligence",
             font_size=24, color=CYAN)

add_divider_line(slide, Inches(1), Inches(3.8), Inches(2))

add_text_box(slide, Inches(1), Inches(4.1), Inches(8), Inches(0.5),
             "Extending AI Lens for Pharmaceutical Manufacturers",
             font_size=18, color=MID_GRAY)

add_text_box(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
             "Poorva Mittal  |  Product Manager, AI Solutions",
             font_size=16, color=MID_GRAY)

add_bottom_bar(slide)

set_notes(slide, """Welcome everyone. Today I'm going to walk you through PharCovAI — a compliance intelligence product that extends AI Lens to solve a critical gap in pharmaceutical patient support operations.

We'll cover the problem space, the product vision, our unique advantages, the MVP scope, and a 90-day pilot plan. I'll close with a live demo of the core product surfaces.

This is an ideation paper and execution plan — I'm looking for your input on feasibility, design, and prioritization.""")


# ── SLIDE 2: Agenda ──────────────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_slide_title(slide, "Agenda")
add_bottom_bar(slide)

items = [
    ("01", "The Problem", "Why pharma compliance is broken today"),
    ("02", "Product Vision", "From 5% to 100% conversation coverage"),
    ("03", "Why the Company", "Our moat and market timing"),
    ("04", "MVP Scope", "V1 features and dependency chain"),
    ("05", "Product Experience", "Core workflow and surfaces"),
    ("06", "Go-to-Market", "Onboarding, pilot plan, and success metrics"),
    ("07", "The Ask", "What we need to proceed"),
    ("08", "Live Demo", "Interactive prototype walkthrough"),
]

y_start = Inches(1.7)
for i, (num, title, desc) in enumerate(items):
    y = y_start + Inches(i * 0.62)
    add_rounded_rect(slide, Inches(1), y, Inches(0.65), Inches(0.42),
                     CYAN, num, font_size=14, font_color=WHITE)
    add_text_box(slide, Inches(1.85), y + Inches(0.02), Inches(3), Inches(0.4),
                 title, font_size=18, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(5), y + Inches(0.05), Inches(6), Inches(0.4),
                 desc, font_size=14, color=MID_GRAY)

# Time indicator
add_rounded_rect(slide, Inches(10.5), Inches(6.5), Inches(2.2), Inches(0.45),
                 LIGHT_BLUE, "25 min talk + 5 min demo", font_size=12)

set_notes(slide, """Here's the roadmap for our session. I've structured this to cover the problem, the product, and the plan in 25 minutes, with 5 minutes at the end for a live demo of the prototype.

I encourage questions throughout — especially from engineering on feasibility and UX on workflow design.""")


# ── SLIDE 3: The Problem — Key Stats ────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 1)
add_slide_title(slide, "The Compliance Blind Spot")
add_bottom_bar(slide)

# Stat cards row
add_stat_card(slide, Inches(0.8), Inches(2.0), "1–5%", "Conversations\nReviewed Today", RED)
add_stat_card(slide, Inches(3.7), Inches(2.0), "5–10%", "Adverse Reactions\nActually Reported", ORANGE)
add_stat_card(slide, Inches(6.6), Inches(2.0), "15 days", "FDA Reporting\nWindow", AMBER)
add_stat_card(slide, Inches(9.5), Inches(2.0), "$100M+", "Potential Fines for\nNon-Compliance", RED)

# Bottom insight
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.8),
             "Every unreviewed conversation is a potential unreported adverse event, "
             "an undetected SOP deviation, or an audit gap waiting to surface.",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Visual: the gap
add_rounded_rect(slide, Inches(2), Inches(5.6), Inches(4), Inches(1.0),
                 BORDER_GRAY, "Manual QA: 1 reviewer per\n500-1,000 calls",
                 font_size=13, font_color=DARK_BLUE, bold=False)

add_arrow(slide, Inches(6.2), Inches(5.85), Inches(0.8), Inches(0.4), CYAN)

add_rounded_rect(slide, Inches(7.2), Inches(5.6), Inches(4), Inches(1.0),
                 CYAN, "Structural bottleneck:\n95-99% of conversations unreviewed",
                 font_size=13, font_color=WHITE, bold=False)

set_notes(slide, """Let me start with the core problem.

Pharma patient support programs generate tens of thousands of interactions monthly. But manual QA teams typically staff just 1 reviewer per 500 to 1,000 calls. The result? Only 1 to 5 percent of conversations are ever reviewed.

Studies show only 5 to 10 percent of adverse drug reactions are reported through spontaneous reporting systems. The FDA requires serious AE reports within 15 calendar days, and penalties for missing that window can reach into the hundreds of millions.

Every unreviewed conversation is a potential unreported adverse event. That's not a quality problem — it's a structural one. And it's getting worse as AI agents enter these conversations without centralized oversight.""")


# ── SLIDE 4: The Failure Chain ───────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 1)
add_slide_title(slide, "How Adverse Events Get Missed",
                "The failure chain from patient mention to audit discovery")
add_bottom_bar(slide)

chain_items = [
    ("1", "Patient mentions\nsymptom on call", RED),
    ("2", "Agent misses\nor logs poorly", ORANGE),
    ("3", "Enters QA\nbacklog", AMBER),
    ("4", "AE goes\nunreported", RED),
    ("5", "Found months\nlater in audit", RGBColor(0x7C, 0x3A, 0xED)),
]

y_center = Inches(3.5)
x_start = Inches(0.8)
box_w = Inches(2.0)
box_h = Inches(1.4)
gap = Inches(0.5)

for i, (num, text, color) in enumerate(chain_items):
    x = x_start + i * (box_w + gap)
    add_rounded_rect(slide, x, y_center, box_w, box_h, color, text, font_size=13)
    add_circle(slide, x + box_w / 2 - Inches(0.2), y_center - Inches(0.55),
               Inches(0.4), color, num, font_size=14)
    if i < len(chain_items) - 1:
        add_arrow(slide, x + box_w + Inches(0.05), y_center + box_h / 2 - Inches(0.15),
                  Inches(0.4), Inches(0.3), BORDER_GRAY)

# AI governance callout
add_rounded_rect(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.2),
                 RGBColor(0xFE, 0xF3, 0xC7),
                 "AI Governance Gap: AI agents handle growing share of patient interactions "
                 "with no centralized monitoring, no audit trail, and no compliance feedback loop",
                 font_size=13, font_color=DARK_BLUE, bold=False)

set_notes(slide, """Here's how the failure chain works in practice.

Step 1: A patient mentions a symptom or side effect during a support call.
Step 2: The agent — whether human or AI — either misses the signal or logs it inconsistently.
Step 3: That interaction enters a QA backlog where it may never be sampled.
Step 4: The adverse event goes unreported, potentially breaching the FDA's 15-day window.
Step 5: The gap is only discovered during an audit — weeks or months later.

And there's a compounding factor: AI agents now handle a growing share of these interactions, but most organizations have zero centralized monitoring of what the AI said, whether it followed the script, or whether it surfaced a potential AE correctly. That's the AI governance vacuum.""")


# ── SLIDE 5: Market Opportunity ──────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 1)
add_slide_title(slide, "Market Opportunity",
                "The pharmacovigilance market is large, growing, and underserved")
add_bottom_bar(slide)

# Market size visual
add_stat_card(slide, Inches(1), Inches(2.2), "$10.4B", "PV Market 2025", CYAN)
add_stat_card(slide, Inches(4), Inches(2.2), "$22.3B", "Projected 2034", LIGHT_BLUE)
add_stat_card(slide, Inches(7), Inches(2.2), "8.9%", "CAGR Growth", GREEN)
add_stat_card(slide, Inches(10), Inches(2.2), "50%", "Faster Ops with\nAI (EVERSANA)", CYAN)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
             "No healthcare-native, AI-powered compliance intelligence platform exists today.",
             font_size=20, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Three-path comparison
paths = [
    ("Build Internally", "12-18 months", "Diverts engineering,\nno shared learnings", RED),
    ("Buy Generic", "Weeks", "Not healthcare-native,\nno PV taxonomy", ORANGE),
    ("Extend with AI Lens", "90-day pilot", "Leverages existing moat\nfrom Day 1", GREEN),
]
for i, (title, timeline, gap, color) in enumerate(paths):
    x = Inches(1.5) + i * Inches(3.8)
    add_rounded_rect(slide, x, Inches(5.8), Inches(3.2), Inches(0.45),
                     color, title, font_size=14)
    add_text_box(slide, x + Inches(0.1), Inches(6.3), Inches(3.0), Inches(0.3),
                 timeline, font_size=13, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(6.6), Inches(3.0), Inches(0.5),
                 gap, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

set_notes(slide, """The pharmacovigilance market is valued at over 10 billion dollars in 2025 and projected to reach 22 billion by 2034 — growing at nearly 9% CAGR. AI-driven automation is a key growth driver. EVERSANA has already demonstrated 50% faster operations with 40% less manual effort using AI-powered PV.

But here's the important insight: no healthcare-native, AI-powered compliance intelligence platform exists today.

Pharma companies face three paths: build internally, which takes 12-18 months and diverts engineering resources. Buy generic tools like Observe.AI or CallMiner, which aren't healthcare-native and lack PV taxonomy. Or extend with the Company — where we can leverage our existing infrastructure for a 90-day pilot.

The extend path is uniquely available through us because our foundational layers already exist.""")


# ── SLIDE 6: Product Vision ─────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 2)
add_slide_title(slide, "From 5% to 100%: The Product Vision")
add_bottom_bar(slide)

# Before / After comparison
headers = ["", "Today", "With PharCovAI"]
rows_data = [
    ["Conversations reviewed", "1-5% (manual QA)", "100% (AI + human oversight)"],
    ["AE detection", "Agent self-reporting", "High-recall AI with evidence"],
    ["Time to flag", "Days to weeks", "Minutes (P95 < 30 min)"],
    ["Audit readiness", "Weeks of reconciliation", "Continuous, export-ready"],
    ["AI agent oversight", "None", "Centralized governance"],
]
add_table(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(3.0),
          len(rows_data) + 1, 3, headers, rows_data,
          col_widths=[Inches(3), Inches(4), Inches(4.3)])

# Design principles as icons
principles = [
    ("Recall > Precision", "Over-flag, never miss"),
    ("Human-in-the-Loop", "AI assists, humans decide"),
    ("Audit Trail by Default", "Immutable logging"),
    ("Channel-Agnostic", "Voice, SMS, email, AI"),
]
y = Inches(5.3)
for i, (title, desc) in enumerate(principles):
    x = Inches(0.8) + i * Inches(3.1)
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(0.45), CYAN,
                     title, font_size=12, font_color=WHITE)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.55), Inches(2.6), Inches(0.4),
                 desc, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

set_notes(slide, """The vision is straightforward: transition pharma customers from reviewing 5% of conversations to evaluating 100% — while maintaining human oversight for all regulatory determinations.

Look at the before and after. Today, AE detection relies on agent self-reporting and random audits, taking days to weeks. With PharCovAI, we flag potential adverse events within minutes with full evidence.

Four design principles govern our trade-offs:
1. Recall over precision — better to over-flag than miss a reportable event. We target 97% recall even if it means more false positives.
2. Human-in-the-loop — AI never makes a final compliance call. Every flagged event routes to a human reviewer.
3. Audit trail by default — every action is logged immutably. This is an architectural constraint, not a feature.
4. Channel-agnostic — we analyze voice, SMS, email, and AI-agent interactions through a unified pipeline.""")


# ── SLIDE 7: Why Now — The Strategic Moat ──────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 3)
add_slide_title(slide, "Why Now — The Strategic Moat",
                "Four converging signals create a narrow window")
add_bottom_bar(slide)

signals = [
    ("AI Agents Enter\nRegulated Conversations",
     "Patient support programs deploying AI agents at scale without compliance oversight",
     CYAN),
    ("LLM Accuracy\nCrosses Threshold",
     "Medical NLP enables 95%+ AE recall without overwhelming reviewers",
     LIGHT_BLUE),
    ("Regulators\nSignal Intent",
     "FDA AI/ML framework, EU AI Act high-risk classification for healthcare AI",
     RGBColor(0x7C, 0x3A, 0xED)),
    ("the Company Has\nthe Infrastructure",
     "Not a cold start — conversation processing at scale already exists",
     CYAN),
]

for i, (title, desc, color) in enumerate(signals):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(2.0)

    add_circle(slide, x + Inches(1.0), y, Inches(0.7), color,
               str(i + 1), font_size=20)

    add_text_box(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.9), Inches(0.8),
                 title, font_size=16, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.1), y + Inches(1.9), Inches(2.9), Inches(1.2),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Moat summary bar
moat_items = [
    ("Millions+", "Conversations"),
    ("Millions", "Calls Processed"),
    ("1,000+", "Therapies"),
    ("Significant\nShare", "Healthcare F50"),
]
y_bar = Inches(5.5)
for i, (num, label) in enumerate(moat_items):
    x = Inches(1.2) + i * Inches(2.9)
    add_text_box(slide, x, y_bar, Inches(2.5), Inches(0.5),
                 num, font_size=28, bold=True, color=CYAN,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y_bar + Inches(0.5), Inches(2.5), Inches(0.4),
                 label, font_size=12, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.6), Inches(11.3), Inches(0.4),
             "A competitor would need 18-24 months to replicate these advantages",
             font_size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

set_notes(slide, """Why should the Company build this, and why now? Four signals are converging:

First, AI agents are entering regulated conversations at scale. Every AI-handled interaction creates a new compliance surface without corresponding oversight.

Second, LLM capabilities have crossed the accuracy threshold. High-recall AE detection at 95% or better is now achievable without overwhelming reviewers.

Third, regulators are signaling intent. The FDA's evolving AI framework and the EU AI Act's high-risk classification tell us companies deploying AI in patient-facing roles will need centralized monitoring.

Fourth — and this is our key advantage — the Company already has the infrastructure. We've processed hundreds of millions of minutes of conversations across hundreds of therapies, serving a significant share of the healthcare Fortune 50. This is an extension, not a greenfield build. A competitor would need 18-24 months to replicate what we already have.""")


# ── SLIDE 8: Target Users ───────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 4)
add_slide_title(slide, "Target Users & Personas")
add_bottom_bar(slide)

personas = [
    ("PV Specialist", "Primary User",
     "Reviews AI-flagged cases\nConfirm/dismiss/escalate AEs\n<4 min per case target",
     "Case Detail View", CYAN),
    ("Operations Lead", "Power User",
     "Monitors SOP adherence\nAgent coaching decisions\nTrend analysis",
     "SOP Dashboard", LIGHT_BLUE),
    ("Compliance VP", "Economic Buyer",
     "Owns regulatory risk\nAudit readiness\nVendor procurement",
     "Executive Reports", RGBColor(0x7C, 0x3A, 0xED)),
    ("IT/Security", "Gatekeeper",
     "HIPAA validation\nSSO/RBAC config\nData pipeline security",
     "Admin Console", MID_GRAY),
]

for i, (name, role, tasks, surface, color) in enumerate(personas):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(1.8)

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(4.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER_GRAY
    card.line.width = Pt(1)

    # Colored header
    add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.2),
                     Inches(2.6), Inches(0.8), color, f"{name}\n{role}",
                     font_size=13)

    # Tasks
    add_text_box(slide, x + Inches(0.25), y + Inches(1.2),
                 Inches(2.4), Inches(2.0),
                 tasks, font_size=12, color=DARK_BLUE, line_spacing=1.6)

    # Surface label
    add_rounded_rect(slide, x + Inches(0.3), y + Inches(3.8),
                     Inches(2.3), Inches(0.4), LIGHT_GRAY,
                     f"Primary: {surface}", font_size=11,
                     font_color=DARK_BLUE, bold=False)

set_notes(slide, """Four personas interact with PharCovAI, but our MVP is laser-focused on the PV Specialist.

The PV Specialist is our primary user. They review AI-flagged cases, make confirm/dismiss/escalate decisions, and we're targeting under 4 minutes per case — down from 15-20 minutes with manual review.

The Operations Lead monitors SOP adherence and makes agent coaching decisions. They come in with V1.5.

The Compliance VP is our economic buyer — they own regulatory risk, audit readiness, and vendor procurement decisions.

IT and Security are gatekeepers for HIPAA validation and SSO configuration.

For the MVP, we focus entirely on the PV Specialist workflow. If we nail their experience, the other personas follow.""")


# ── SLIDE 9: MVP Features ───────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 4)
add_slide_title(slide, "MVP Feature Prioritization",
                "V1: Complete the AE detection workflow before expanding")
add_bottom_bar(slide)

# Dependency chain visual
features = [
    ("F1", "Voice Call\nIngestion", "No data = nothing\nto analyze"),
    ("F2", "AE Detection\nEngine", "FDA 15-day mandate\nmakes this top risk"),
    ("F3", "Case Queue", "100% coverage needs\nstructured triage"),
    ("F4", "Case Detail\nView", "Evidence + MedWatch\nin one surface"),
    ("F5", "Audit-Ready\nExport", "Prerequisite for\ncompliance sign-off"),
    ("F6", "Studio\nIntegration", "Closes detection\nto action loop"),
]

y_top = Inches(2.0)
box_w = Inches(1.7)
box_h = Inches(1.2)

for i, (fid, name, why) in enumerate(features):
    x = Inches(0.5) + i * Inches(2.1)

    add_rounded_rect(slide, x, y_top, box_w, Inches(0.35), CYAN,
                     fid, font_size=12)
    add_rounded_rect(slide, x, y_top + Inches(0.4), box_w, box_h,
                     LIGHT_GRAY, name, font_size=12, font_color=DARK_BLUE,
                     bold=True)
    add_text_box(slide, x + Inches(0.05), y_top + Inches(1.7),
                 box_w - Inches(0.1), Inches(0.8),
                 why, font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(features) - 1:
        add_arrow(slide, x + box_w + Inches(0.05),
                  y_top + Inches(0.7), Inches(0.3), Inches(0.25), CYAN)

# Future versions
add_divider_line(slide, Inches(0.5), Inches(4.8), Inches(12.3))

versions = [
    ("V1.5", "SOP Scoring  |  SMS & Email  |  PV System Integration",
     LIGHT_BLUE),
    ("V2", "Fax Intake  |  Multilingual Support", MID_GRAY),
    ("Deferred", "Real-Time Monitoring  |  Automated FDA Submission",
     BORDER_GRAY),
]
for i, (ver, items, color) in enumerate(versions):
    y = Inches(5.1) + i * Inches(0.7)
    add_rounded_rect(slide, Inches(0.8), y, Inches(1.2), Inches(0.45),
                     color, ver, font_size=12)
    add_text_box(slide, Inches(2.3), y + Inches(0.05), Inches(10), Inches(0.4),
                 items, font_size=14, color=DARK_BLUE)

set_notes(slide, """The MVP is scoped around one principle: deliver a complete AE detection workflow before expanding.

F1 through F6 form a dependency chain — each feature is useless without the one before it.

You need ingestion before detection, detection before a reviewer workflow, reviewer decisions before audit readiness, and a case to act on before remediation.

F1: Voice call ingestion — this is the dominant channel.
F2: AE detection engine — high-recall classification, the core of the product.
F3: Case queue — severity-prioritized triage for reviewers.
F4: Case detail view — transcript evidence, audio, and pre-populated MedWatch fields.
F5: Audit-ready export — regulatory export packages.
F6: Studio integration — closes the detection-to-action loop.

Everything outside this chain — SOP scoring, multi-channel support, PV system integration — ships in V1.5 after the pilot validates the foundation.""")


# ── SLIDE 10: Product Experience — Workflow ──────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 5)
add_slide_title(slide, "Product Experience: Core Workflow",
                "From AI flag to regulatory decision in under 4 minutes")
add_bottom_bar(slide)

# Workflow steps
steps = [
    ("Conversation\nIngested", "Voice call enters\nprocessing pipeline", CYAN),
    ("AI Analysis", "AE detection, severity\nscoring, MedWatch extract", LIGHT_BLUE),
    ("Triage Queue", "Reviewer sees\nprioritized cases", CYAN),
    ("Case Review", "Transcript + audio +\nevidence in one view", LIGHT_BLUE),
    ("Decision", "Confirm / Dismiss /\nEscalate", GREEN),
    ("Audit Log", "Immutable record of\nevery action", RGBColor(0x7C, 0x3A, 0xED)),
]

y_flow = Inches(2.3)
step_w = Inches(1.7)
step_h = Inches(1.1)

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.4) + i * Inches(2.1)

    # Step circle
    add_circle(slide, x + step_w / 2 - Inches(0.25), y_flow - Inches(0.1),
               Inches(0.5), color, str(i + 1), font_size=16)

    # Step box
    add_rounded_rect(slide, x, y_flow + Inches(0.6), step_w, step_h,
                     color, title, font_size=13)

    # Description
    add_text_box(slide, x + Inches(0.05), y_flow + step_h + Inches(0.7),
                 step_w - Inches(0.1), Inches(0.8),
                 desc, font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(steps) - 1:
        add_arrow(slide, x + step_w + Inches(0.05), y_flow + Inches(0.85),
                  Inches(0.3), Inches(0.25), BORDER_GRAY)

# Two product surfaces
add_divider_line(slide, Inches(0.5), Inches(5.0), Inches(12.3))

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5), Inches(0.4),
             "Core Product Surfaces", font_size=18, bold=True, color=DARK_BLUE)

surfaces = [
    ("Triage Queue", "Severity-prioritized case list with reviewer\nassignment, confidence scoring, and SLA countdown"),
    ("Case Detail View", "Highlighted transcript with audio, pre-populated\nMedWatch fields, and confirm/dismiss/escalate actions"),
]
for i, (name, desc) in enumerate(surfaces):
    x = Inches(0.8) + i * Inches(6.2)
    y = Inches(5.7)
    add_rounded_rect(slide, x, y, Inches(2), Inches(0.4), CYAN,
                     name, font_size=13)
    add_text_box(slide, x + Inches(2.2), y + Inches(0.0), Inches(3.8), Inches(0.6),
                 desc, font_size=11, color=MID_GRAY)

set_notes(slide, """Here's the core reviewer workflow.

Step 1: A conversation is ingested — for MVP, voice calls.
Step 2: AI analysis runs AE detection, severity scoring, and MedWatch field extraction.
Step 3: Flagged cases appear in the Triage Queue, prioritized by severity.
Step 4: The reviewer clicks into a case and sees transcript, audio evidence, and AI reasoning in one view.
Step 5: They make a decision — Confirm, Dismiss, or Escalate.
Step 6: Every action is immutably logged in the audit trail.

There are two core product surfaces: the Triage Queue for deciding what to review next, and the Case Detail View for making the actual regulatory determination. I'll show both of these in the demo.""")


# ── SLIDE 11: Onboarding ────────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 6)
add_slide_title(slide, "From Contract to First Case Review",
                "Why pharma onboarding is different — and how we handle it")
add_bottom_bar(slide)

phases = [
    ("Pre-Contract", "2-4 wk", "BAA/DPA, therapy\nscoping, success criteria", "Signed BAA/DPA\n+ agreed scope"),
    ("Technical\nSetup", "1-2 wk", "SSO, RBAC,\npipeline validation", "Data flowing\nthrough pipeline"),
    ("Calibration", "2-3 wk", "Model tuning on\ncustomer data", "Recall >= 95%\non customer data"),
    ("Soft Launch", "1-2 wk", "Production dry-run,\nreviewer training", "Reviewers working\nindependently"),
    ("Steady\nState", "Ongoing", "Governance cadence,\nmodel reviews", "Continuous\nimprovement"),
]

y_phase = Inches(2.0)
phase_w = Inches(2.2)

for i, (name, duration, activities, gate) in enumerate(phases):
    x = Inches(0.5) + i * Inches(2.5)

    # Phase header
    color = CYAN if i < 4 else GREEN
    add_rounded_rect(slide, x, y_phase, phase_w, Inches(0.7),
                     color, f"{name}\n({duration})", font_size=11)

    # Activities
    add_text_box(slide, x + Inches(0.1), y_phase + Inches(0.85),
                 phase_w - Inches(0.2), Inches(1.0),
                 activities, font_size=11, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Gate
    if i < 4:
        add_text_box(slide, x + Inches(0.1), y_phase + Inches(2.0),
                     phase_w - Inches(0.2), Inches(0.3),
                     "Gate:", font_size=10, bold=True, color=CYAN,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y_phase + Inches(2.3),
                     phase_w - Inches(0.2), Inches(0.7),
                     gate, font_size=10, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Arrow between phases
    if i < len(phases) - 1:
        add_arrow(slide, x + phase_w + Inches(0.02),
                  y_phase + Inches(0.25), Inches(0.25), Inches(0.2), BORDER_GRAY)

# First value moment callout
add_rounded_rect(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(1.2),
                 RGBColor(0xCC, 0xFB, 0xF1),
                 'First Value Moment: "Yes, this is a real signal I would have missed"\n'
                 'Target: within 5-7 weeks of contract signature',
                 font_size=14, font_color=DARK_BLUE, bold=False)

set_notes(slide, """Pharma onboarding is different from standard SaaS. Before a single conversation flows through the system, three things must be true: legal clearance with executed BAA and DPA agreements, data trust demonstrated on the customer's own conversation data, and workflow fit that matches how reviewers already work.

We've structured onboarding into five gated phases, each with clear criteria to advance.

The most critical moment is what I call the First Value Moment — the first time a reviewer opens a real case flagged by PharCovAI on their own data and says 'yes, this is a real signal I would have missed.' If that happens, adoption follows. We target this within 5-7 weeks of contract signature.""")


# ── SLIDE 12: Pilot Plan ────────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 6)
add_slide_title(slide, "90-Day Pilot Plan")
add_bottom_bar(slide)

# Pilot design stats
pilot_stats = [
    ("4 weeks", "Live Pilot"),
    ("2 partners", "Existing Customers"),
    ("~3,000/wk", "Conversations Each"),
    ("2 analysts", "Per Partner"),
]
for i, (num, label) in enumerate(pilot_stats):
    x = Inches(0.8) + i * Inches(3.1)
    add_stat_card(slide, x, Inches(1.7), num, label, CYAN)

# Gantt-style timeline
tracks = [
    ("Discovery & UX", "Weeks 1-4", 0, 4, LIGHT_BLUE),
    ("Engineering", "Weeks 1-9", 0, 9, CYAN),
    ("Legal & Compliance", "Weeks 4-5", 3, 2, RGBColor(0x7C, 0x3A, 0xED)),
    ("Pilot Operations", "Weeks 7-12", 6, 6, GREEN),
]

gantt_left = Inches(1.5)
gantt_top = Inches(4.3)
week_width = Inches(0.85)
bar_height = Inches(0.4)
bar_gap = Inches(0.55)

# Week labels
for w in range(12):
    add_text_box(slide, gantt_left + w * week_width, gantt_top - Inches(0.35),
                 week_width, Inches(0.3),
                 f"W{w+1}", font_size=9, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

for i, (name, period, start, length, color) in enumerate(tracks):
    y = gantt_top + i * bar_gap

    add_text_box(slide, Inches(0.2), y + Inches(0.02), Inches(1.3), Inches(0.35),
                 name, font_size=11, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.RIGHT)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        gantt_left + start * week_width, y,
        length * week_width, bar_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = period
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

set_notes(slide, """The 90-day pilot is designed to validate the product with two existing the Company pharma customers. Each partner contributes one therapy program, voice-only for MVP, processing roughly 3,000 conversations per week with 2 PV analysts per partner.

Four tracks run in parallel:
- Discovery and UX runs weeks 1 through 4
- Engineering runs weeks 1 through 9 — the longest track
- Legal and compliance is a focused sprint in weeks 4-5 to fast-track BAA and DPA
- Pilot operations run weeks 7 through 12 — this is live production with real conversations

The parallel structure means we're not blocked sequentially. Engineering starts building while UX is still validating designs.""")


# ── SLIDE 13: Success Metrics ────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 6)
add_slide_title(slide, "How We'll Know It's Working",
                "Three-tier scorecard — safety gates, quality targets, business signals")
add_bottom_bar(slide)

# Three tiers
tiers = [
    ("SAFETY", "Non-Negotiable Gates", RED, [
        ("AE Recall", ">= 97%"),
        ("Reporting", "100% in FDA windows"),
        ("Coverage", ">= 99.5%"),
        ("Misconfig", "Zero incidents"),
    ]),
    ("QUALITY", "Optimization Targets", AMBER, [
        ("Precision", ">= 70%"),
        ("Reviewer Agreement", ">= 75%"),
        ("False Positive Rate", "< 20%"),
        ("Time-to-Flag", "< 30 min P95"),
    ]),
    ("BUSINESS", "Commercial Signals", GREEN, [
        ("Time-to-Decision", "< 4 min median"),
        ("Reviewer NPS", ">= +20"),
        ("Conversion", "Both partners by W14"),
        ("Expansion", "Additional programs"),
    ]),
]

for i, (tier_name, subtitle, color, metrics) in enumerate(tiers):
    x = Inches(0.5) + i * Inches(4.2)
    y = Inches(1.8)

    # Tier header
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(0.7), color,
                     f"{tier_name}\n{subtitle}", font_size=13)

    # Metrics
    for j, (metric, target) in enumerate(metrics):
        my = y + Inches(0.9) + j * Inches(0.9)

        # Metric card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), my,
            Inches(3.6), Inches(0.75)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Pt(1)

        add_text_box(slide, x + Inches(0.25), my + Inches(0.05),
                     Inches(2.0), Inches(0.35),
                     metric, font_size=12, color=DARK_BLUE)
        add_text_box(slide, x + Inches(0.25), my + Inches(0.38),
                     Inches(2.0), Inches(0.35),
                     target, font_size=14, bold=True, color=color)

# Note at bottom
add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             "Safety metrics must ALL pass. Quality misses trigger iteration, not termination. "
             "Business metrics inform commercialization.",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

set_notes(slide, """Our success metrics are organized into three tiers, and they are not equal.

Tier 1 — Safety: These are non-negotiable gates. If any safety metric fails, the pilot stops. We need at least 97% AE recall, 100% reporting within FDA windows, 99.5% conversation coverage, and zero misconfiguration incidents.

Tier 2 — Quality: These are optimization targets. If we hit 60% precision instead of 70%, we iterate on the model — we don't terminate the pilot. Same for reviewer agreement, false positive rate, and time-to-flag.

Tier 3 — Business: These are commercial signals measured at pilot end. Reviewer time-to-decision under 4 minutes, NPS of +20 or better, and both partners committing to paid conversion.

The critical distinction: safety failures stop the pilot. Quality misses trigger engineering iteration. Business metrics inform the commercialization decision.""")


# ── SLIDE 14: The Week 12 Decision ───────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 6)
add_slide_title(slide, "The Week 12 Decision",
                "Go / Conditional Go / No-Go framework")
add_bottom_bar(slide)

decisions = [
    ("GO", GREEN, [
        "AE recall >= 97%",
        "No production AE misses",
        "Time-to-decision < 4 min",
        "Both partners commit",
    ], "Convert to paid, begin V1.5,\npublish internal case study"),

    ("CONDITIONAL\nGO", AMBER, [
        "Most criteria met",
        "1-2 non-safety metrics short",
        "Max one retry cycle",
        "Named fixes + owners",
    ], "4-week remediation,\nre-evaluate at Week 16"),

    ("NO-GO", RED, [
        "Recall below 93%",
        "Confirmed production miss",
        "FP rate exceeds 30%",
        "Compliance blocked",
    ], "Post-mortem, preserve data,\nre-evaluate in 6 months"),
]

for i, (title, color, criteria, next_steps) in enumerate(decisions):
    x = Inches(0.6) + i * Inches(4.2)
    y = Inches(1.8)

    # Header
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(0.8), color,
                     title, font_size=20)

    # Criteria
    for j, criterion in enumerate(criteria):
        cy = y + Inches(1.1) + j * Inches(0.55)
        # Bullet indicator
        add_circle(slide, x + Inches(0.2), cy + Inches(0.08),
                   Inches(0.2), color, "", font_size=8)
        add_text_box(slide, x + Inches(0.5), cy, Inches(3.2), Inches(0.45),
                     criterion, font_size=13, color=DARK_BLUE)

    # Next steps
    add_text_box(slide, x + Inches(0.1), y + Inches(3.6), Inches(3.6), Inches(0.3),
                 "Next Steps:", font_size=12, bold=True, color=color)
    add_text_box(slide, x + Inches(0.1), y + Inches(3.9), Inches(3.6), Inches(0.7),
                 next_steps, font_size=12, color=MID_GRAY)

# Rule callout
add_rounded_rect(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.6),
                 RGBColor(0xFE, 0xF3, 0xC7),
                 "Rule: Safety metrics must ALL pass for any Go outcome. "
                 "No Conditional Go for safety failures.",
                 font_size=13, font_color=DARK_BLUE, bold=False)

set_notes(slide, """At the end of Week 12, we hold a Go/No-Go review with the Company leadership and pilot partner compliance leads. Three outcomes:

GO — if AE recall is at or above 97%, no production misses, time-to-decision under 4 minutes, and both partners commit. We convert to paid, begin V1.5 planning, and publish an internal case study.

CONDITIONAL GO — if most criteria are met but 1-2 non-safety metrics fall short. For example, recall hits 97% but precision is at 60% instead of 70%. We get a 4-week remediation cycle with specific named fixes, then re-evaluate at Week 16 as a binary Go/No-Go.

NO-GO — if recall drops below 93%, there's a confirmed production miss, false positives exceed 30%, or compliance counsel blocks us. We do a post-mortem, preserve the data and model weights, and re-evaluate in 6 months.

The critical rule: safety metrics must ALL pass. There's no Conditional Go for safety failures.""")


# ── SLIDE 15: Post-Pilot Roadmap ─────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 7)
add_slide_title(slide, "The Post-Pilot Roadmap",
                "From pilot to system of record in 24 months")
add_bottom_bar(slide)

quarters = [
    ("Q1", "Months 1-3", "PILOT", "Build MVP\nOnboard 2 partners\nValidate AE detection", CYAN),
    ("Q2", "Months 4-6", "COMMERCIALIZE", "Convert to paid\n2-3 new customers\nBegin V1.5", LIGHT_BLUE),
    ("Q3", "Months 7-9", "EXPAND", "Launch V1.5\nSOP + multi-channel\nExpand within customers", RGBColor(0x7C, 0x3A, 0xED)),
    ("Q4", "Months 10-12", "SCALE", "Launch V2\nFax + multilingual\nSales playbook", GREEN),
]

y_road = Inches(2.2)
for i, (q, period, phase, details, color) in enumerate(quarters):
    x = Inches(0.6) + i * Inches(3.15)

    # Phase card
    add_rounded_rect(slide, x, y_road, Inches(2.9), Inches(0.55), color,
                     f"{q}: {phase}", font_size=16)

    add_text_box(slide, x + Inches(0.1), y_road + Inches(0.6),
                 Inches(2.7), Inches(0.3),
                 period, font_size=12, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.15), y_road + Inches(1.0),
                 Inches(2.6), Inches(1.5),
                 details, font_size=13, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    if i < len(quarters) - 1:
        add_arrow(slide, x + Inches(2.95), y_road + Inches(0.15),
                  Inches(0.15), Inches(0.25), BORDER_GRAY)

# North star
add_divider_line(slide, Inches(0.6), Inches(5.0), Inches(12.1))

add_text_box(slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.8),
             "24-Month North Star: PharCovAI becomes the system of record for "
             "AI-enabled pharmacovigilance — analyzing millions of conversations monthly, "
             "across all channels and languages, with regulatory acceptance.",
             font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Strategic argument
points = [
    "Market is spending but underserved — $10B+ with no healthcare-native AI platform",
    "the Company is uniquely positioned — 18-24 month head start over competitors",
    "Window is time-bounded — first mover in AI governance becomes the default",
]
for i, point in enumerate(points):
    y = Inches(6.0) + i * Inches(0.4)
    add_circle(slide, Inches(1.5), y + Inches(0.05), Inches(0.2), CYAN)
    add_text_box(slide, Inches(1.9), y, Inches(10), Inches(0.35),
                 point, font_size=12, color=MID_GRAY)

set_notes(slide, """Looking beyond the pilot, here's the quarterly roadmap.

Q1 is the pilot we've been discussing — build MVP, onboard 2 partners, validate AE detection.

Q2 is commercialization — convert partners to paid, close 2-3 new customers, and begin V1.5 development.

Q3 is expansion — launch V1.5 with SOP scoring and multi-channel support, and grow within existing customers.

Q4 is scale — launch V2 with fax and multilingual support, build the sales playbook, and establish a customer advisory board.

The 24-month north star: PharCovAI becomes the system of record for AI-enabled pharmacovigilance across the Company's pharma customer base.

Three things make this compelling: the market is large and underserved, we have an 18-24 month head start, and the window is time-bounded — the company that establishes the governance layer during this AI adoption wave becomes the default.""")


# ── SLIDE 16: The Ask ────────────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)
add_section_number(slide, 7)
add_slide_title(slide, "The Ask")
add_bottom_bar(slide)

asks = [
    ("Engineering\nAllocation", "Dedicated team of 4-6 engineers\nfor 12 weeks", CYAN),
    ("Data Science\nSupport", "1-2 ML engineers for AE model\ndevelopment and calibration", LIGHT_BLUE),
    ("Pilot Partner\nCommitment", "Executive sponsorship to engage\n2 existing pharma customers", RGBColor(0x7C, 0x3A, 0xED)),
    ("Legal\nFast-Track", "Priority BAA/DPA review\nWeek 1 start, not queued", ORANGE),
    ("Go/No-Go\nAuthority", "Named decision-maker for\nthe Week 12 review", GREEN),
]

for i, (title, desc, color) in enumerate(asks):
    x = Inches(0.4) + i * Inches(2.55)
    y = Inches(2.2)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(3.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER_GRAY
    card.line.width = Pt(1)

    add_circle(slide, x + Inches(0.85), y + Inches(0.3), Inches(0.65),
               color, str(i + 1), font_size=22)

    add_text_box(slide, x + Inches(0.15), y + Inches(1.2), Inches(2.05), Inches(0.8),
                 title, font_size=15, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.15), y + Inches(2.2), Inches(2.05), Inches(1.0),
                 desc, font_size=12, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Cost of inaction
add_rounded_rect(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.9),
                 RGBColor(0xFE, 0xF3, 0xC7),
                 "Cost of Inaction: Customers build internally, generic platforms add surface-level features,\n"
                 "and the AI governance window closes — making this opportunity significantly harder in 12 months",
                 font_size=13, font_color=DARK_BLUE, bold=False)

set_notes(slide, """To proceed, we need five things:

1. Engineering allocation — a dedicated team of 4 to 6 engineers for 12 weeks. This is not a side project.

2. Data science support — 1 to 2 ML engineers for AE model development and calibration on customer data.

3. Pilot partner commitment — executive sponsorship to engage 2 existing pharma customers. These should be customers with active patient support programs.

4. Legal fast-track — priority BAA and DPA review starting Week 1, not queued behind other deals. Legal clearance is on the critical path.

5. Go/No-Go authority — a named decision-maker for the Week 12 review. We need clear accountability for the outcome.

And I want to highlight the cost of inaction. If we don't build this, our customers will build it internally. Generic platforms will add surface-level compliance features. And the AI governance window will close — making this opportunity significantly harder in 12 months.""")


# ── SLIDE 17: Demo Transition ────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)

# Large accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.15))
bar.fill.solid()
bar.fill.fore_color.rgb = CYAN
bar.line.fill.background()

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
             "Live Demo", font_size=48, bold=True, color=DARK_BLUE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
             "Interactive Prototype Walkthrough", font_size=24, color=CYAN,
             alignment=PP_ALIGN.CENTER)

add_divider_line(slide, Inches(5.5), Inches(4.2), Inches(2.3))

# Demo flow
demo_steps = [
    ("1. Lens Dashboard", "Product entry point and navigation"),
    ("2. Triage Queue", "AI-flagged cases, severity sorting, SLA tracking"),
    ("3. Case Detail", "Transcript evidence, MedWatch fields, regulatory decision"),
]
for i, (step, desc) in enumerate(demo_steps):
    y = Inches(4.6) + i * Inches(0.65)
    add_rounded_rect(slide, Inches(3.5), y, Inches(2.5), Inches(0.45),
                     CYAN, step, font_size=13)
    add_text_box(slide, Inches(6.2), y + Inches(0.05), Inches(5), Inches(0.4),
                 desc, font_size=13, color=MID_GRAY)

add_bottom_bar(slide)

set_notes(slide, """Now let me switch to the live demo.

I'll walk through the three product surfaces:

First, the Lens Dashboard — this is the entry point into PharCovAI, integrated into the existing AI Lens navigation.

Second, the Triage Queue — where reviewers see AI-flagged adverse event cases sorted by severity, with confidence scores and SLA countdown timers.

Third, the Case Detail View — where the actual regulatory determination happens. I'll show you the highlighted transcript with the AI-detected adverse event, the pre-populated MedWatch fields, and the confirm/dismiss/escalate workflow.

Let me open the demo...""")


# ── SLIDE 18: Thank You ─────────────────────────────────────
slide = add_blank_slide()
set_slide_bg(slide, WHITE)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.15))
bar.fill.solid()
bar.fill.fore_color.rgb = CYAN
bar.line.fill.background()

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "Thank You", font_size=48, bold=True, color=DARK_BLUE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             "Questions & Discussion", font_size=24, color=CYAN,
             alignment=PP_ALIGN.CENTER)

add_divider_line(slide, Inches(5.5), Inches(4.0), Inches(2.3))

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Poorva Mittal  |  Product Manager, AI Solutions",
             font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Key takeaway cards
takeaways = [
    ("$10B+ market", "No HC-native AI\ncompliance platform"),
    ("90-day pilot", "2 partners,\nvalidated workflow"),
    ("18-24 mo advantage", "the Company moat is\nreal and compounding"),
]
for i, (title, desc) in enumerate(takeaways):
    x = Inches(2) + i * Inches(3.3)
    y = Inches(5.3)
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(0.5), CYAN,
                     title, font_size=14)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.6), Inches(0.6),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

set_notes(slide, """Thank you for your time. Let me leave you with three takeaways:

First, this is a 10 billion dollar market with no healthcare-native AI compliance platform. The opportunity is real and growing at 9% annually.

Second, we can validate the product in 90 days with 2 existing partners. The pilot plan is concrete, the metrics are defined, and the decision framework is clear.

Third, the Company has an 18-24 month head start. Our data moat, infrastructure, customer trust, and domain expertise compound in ways that are extremely difficult to replicate.

I'm looking for your input on three things: engineering feasibility, UX workflow design, and whether this strategic direction resonates. Happy to take questions.""")


# ── Save ─────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__),
                           "PharCovAI-Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
